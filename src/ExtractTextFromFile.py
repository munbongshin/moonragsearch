import os
import io
import re,sys
import tempfile
import time
import PyPDF2
import docx
import logging
from openpyxl import load_workbook
from pptx import Presentation
import pythoncom
from bs4 import BeautifulSoup
import markdown
import pandas as pd
from typing import Union, Optional, List, IO
from langchain.docstore.document import Document
from dotenv import load_dotenv
from pathlib import Path

import platform

if platform.system() != 'Windows':
    print("이 스크립트는 Windows에서만 실행할 수 있습니다.")

if platform.system() == 'Windows':
    import win32com.client as win32    
    import win32gui
else:
    from hwp5.xmlmodel import Hwp5File
    pass



class ExtractTextFromFile:
    def __init__(self):
        project_root = Path(__file__).parent.parent
        # .env 파일의 경로 설정
        env_path = project_root / '.env'
        load_dotenv(dotenv_path=env_path)
        self.chunk_size = int(os.environ.get("CHUNK_SIZE"))
        self.chunk_overlap = int(os.environ.get("CHUNK_OVERLAP"))

    def extract_text_from_file(self, file, file_name):
        if isinstance(file, str):
            file_path = file
            file_name = os.path.basename(file_path)
            _, ext = os.path.splitext(file_name)
        elif hasattr(file, 'name'):
            file_path = file
            file_name = file_name
            _, ext = os.path.splitext(file_name)
        else:
            raise ValueError("Invalid file input. Must be a file path or a file object.")
        
        ext = ext.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf_pages(file_path, file_name)            
        elif ext in ['.docx', '.doc']:
            return self.extract_text_from_docx_pages(file_path)            
        elif ext in ['.xlsx', '.xls']:
            return self.extract_text_from_xlsx_pages(file_path)            
        elif ext in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx_pages(file_path)
        elif ext in ['.hwp', '.hwpx']:
            return self.hwp_to_pdf(file_path, file_name)
        elif ext in ['.txt']:
            return self.extract_text_from_txt(file_path)
        elif ext in ['.md']:
            return self.extract_text_from_markdown(file_path)
        elif ext in ['.htm', '.html']:
            return self.extract_text_from_html(file_path)      
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def extract_text_from_pdf_pages(self, file: Union[str, IO],file_name) -> List[Document]:
        """
        Extract text from PDF pages.

        Args:
            file (Union[str, IO]): Path to the PDF file or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from PDF pages
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        # 파일 객체가 읽기 모드로 열려 있는지 확인
        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None
        
        try:
            if isinstance(file, str):
                file_path = file
                file_name = file_name
                file_obj = open(file, 'rb')
            else:
                file_obj = file
                file_name = file_name
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            pdf_reader = PyPDF2.PdfReader(file_obj)

            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                text = re.sub(r'[^a-zA-Z0-9가-힣\s\n]', '', text)
                text = re.sub(r'(· +)+', '· ', text)
                metadata = {
                    "source": file_path,
                    "file_name": file_name,
                    "page": page_num
                }
                doc = Document(page_content=text, source=file_name, metadata=metadata)
                documents.append(doc)

            return documents

        except Exception as e:
            print(f"Error processing PDF file: {str(e)}")
            return []

        finally:
            # 우리가 파일을 열었다면 닫아줍니다
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.BytesIO):
                file_obj.close() 
        

    def extract_text_from_docx(self, file):
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    
    def extract_text_from_docx_pages(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from Word document pages.

        Args:
            file (Union[str, IO]): Path to the Word document or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from Word document pages
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        # 파일 객체가 읽기 모드로 열려 있는지 확인
        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        temp_file = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                doc = docx.Document(file)
            else:
                # 임시 파일 생성
                temp_file = io.BytesIO(file.read())
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우
                doc = docx.Document(temp_file)

            current_page = ""
            page_number = 1

            for para in doc.paragraphs:
                current_page += para.text + "\n"
                current_page = re.sub(r'[^a-zA-Z0-9가-힣\s\n]', '', current_page)
                current_page = re.sub(r'(· +)+', '· ', current_page)
                # 임의로 1000자마다 새 페이지로 간주 (실제 페이지 구분은 더 복잡할 수 있음)
                if len(current_page) > 1000:
                    metadata = {
                        "source": file_path,
                        "file_name": file_name,
                        "page": page_number
                    }
                    document = Document(page_content=current_page.strip(), metadata=metadata)
                    documents.append(document)
                    current_page = ""
                    page_number += 1

            # 마지막 페이지 추가
            if current_page:
                metadata = {
                    "source": file_path,
                    "file_name": file_name,
                    "page": page_number
                }
                document = Document(page_content=current_page.strip(), metadata=metadata)
                documents.append(document)

            return documents

        except Exception as e:
            print(f"Error processing Word document: {str(e)}")
            return []

        finally:
            if temp_file:
                temp_file.close()   

        
    
    def extract_text_from_xlsx(self, file_path):
        wb = load_workbook(file_path)
        text = ""
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"                
        return text
    
    def extract_text_from_xlsx_pages(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from Excel file pages.

        Args:
            file (Union[str, IO]): Path to the Excel file or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from Excel sheets
        """
        documents = []
        temp_file = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file_path)
            else:
                # 임시 파일 생성
                temp_file = io.BytesIO(file.read())
                file_name = self.get_file_name(file)
                file_path = file_name  # 파일 객체의 이름을 file_path로 사용

            # ExcelFile 객체를 사용하여 모든 시트를 한 번에 읽습니다
            with pd.ExcelFile(temp_file if temp_file else file_path) as xls:
                for sheet_name in xls.sheet_names:
                    # 각 시트를 DataFrame으로 읽습니다
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    
                    # DataFrame을 문자열로 변환합니다
                    buffer = io.StringIO()
                    df.to_csv(buffer, index=False)
                    sheet_content = buffer.getvalue()
                    #sheet_content = re.sub(r'[^a-zA-Z0-9가-힣\s\n]', '', sheet_content)
                    #sheet_content = re.sub(r'(· +)+', '· ', sheet_content)

                    # Document 객체를 생성합니다
                    metadata = {
                        "source": file_path,
                        "file_name": file_name,
                        "page": sheet_name
                    }
                    doc = Document(page_content=sheet_content, metadata=metadata)
                    documents.append(doc)

            return documents

        except Exception as e:
            print(f"Error processing Excel file: {str(e)}")
            return []

        finally:
            if temp_file:
                temp_file.close()
    

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
            text = re.sub(r'[^a-zA-Z0-9가-힣\s\n]', '', text)
            text = re.sub(r'(· +)+', '· ', text)
        return text
    
    def extract_text_from_pptx_pages(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from PowerPoint file pages.

        Args:
            file (Union[str, IO]): Path to the PowerPoint file or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from PowerPoint slides
        """
        documents = []
        temp_file = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file_path)
                prs = Presentation(file_path)
            else:
                # 임시 파일 생성
                temp_file = io.BytesIO(file.read())
                file_name = self.get_file_name(file)
                file_path = file_name  # 파일 객체의 이름을 file_path로 사용
                prs = Presentation(temp_file)

            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text += shape.text + "\n"
                
                # 특수 문자 제거
                slide_text = re.sub(r'[^a-zA-Z0-9가-힣\s\n]', '', slide_text)
                slide_text = re.sub(r'(· +)+', '· ', slide_text)
                # Document 객체를 생성합니다
                metadata = {
                    "source": file_path,
                    "file_name": file_name,
                    "page": slide_number
                }
                doc = Document(page_content=slide_text, metadata=metadata)
                documents.append(doc)

            return documents

        except Exception as e:
            print(f"Error processing PowerPoint file: {str(e)}")
            return []

        finally:
            if temp_file:
                temp_file.close()    
    def chunk_text(self, text, chunk_size):
        return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

    def extract_text_from_file(self, file, file_name):
        if isinstance(file, str):  # file이 문자열(파일 경로)인 경우
            file_path = file
            file_name = os.path.basename(file_path)
            _, ext = os.path.splitext(file_name)
        elif hasattr(file, 'name'):  # file이 파일 객체인 경우
            file_path = file
            file_name = file_name
            _, ext = os.path.splitext(file_name)
        else:
            raise ValueError("Invalid file input. Must be a file path or a file object.")
        
        ext = ext.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf_pages(file_path, file_name)            
        elif ext in ['.docx', '.doc']:
            return self.extract_text_from_docx_pages(file_path)            
        elif ext in ['.xlsx', '.xls']:
            return self.extract_text_from_xlsx_pages(file_path)            
        elif ext in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx_pages(file_path)
        elif ext in ['.hwp', '.hwpx']:
            return self.hwp_to_pdf(file_path,file_name)
        elif ext in ['.txt']:
            return self.extract_text_from_txt(file_path)
        elif ext in ['.md']:
            return self.extract_text_from_markdown(file_path)
        elif ext in ['.htm', '.html']:
            return self.extract_text_from_html(file_path)      
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        
    def hwp_to_pdf(self, hwp_file: Union[str, IO],file_name) -> Optional[List[Document]]:
        """
        Convert HWP file to PDF and extract text from the converted PDF.
        
        Args:
            hwp_file (Union[str, IO]): Path to the HWP file or a file-like object
        Returns:
            Optional[List[Document]]: List of Document objects with extracted text and PDF path as source, or None if conversion fails
        """
        
        try:
            pythoncom.CoInitialize()
            hwp = win32.gencache.EnsureDispatch("HWPFrame.HwpObject")
            hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
        except Exception:
            print("한글(HWP)이 설치되어 있지 않거나 COM 객체를 생성할 수 없습니다.")
            return None   
        
        temp_pdf_path = None
        
        try:
            if isinstance(hwp_file, str):
                hwp_path = hwp_file
                base_name = os.path.splitext(os.path.basename(hwp_file))[0]
            else:
                # 파일 객체에서 임시 HWP 파일 생성
                base_name =  os.path.splitext(file_name)[0]
                temp_dir = tempfile.gettempdir()
                temp_hwp_path = os.path.join(temp_dir, f"{base_name}.hwp")
                with open(temp_hwp_path, 'wb') as temp_file:
                    temp_file.write(hwp_file.read())
                hwp_path = temp_hwp_path

            # 임시 PDF 파일 경로 설정
            temp_pdf_path = os.path.join(tempfile.gettempdir(), f"{base_name}.pdf")
        
            
            hwp.Open(hwp_path)
            hwp.SaveAs(temp_pdf_path, "PDF")
            
            time.sleep(0.5)
            self.handle_completion_dialog()
            
            time.sleep(0.5)
            self.close_hwp_popups()
            print(f"pdf fime name : {temp_pdf_path}")
            
            documents = self.extract_text_from_pdf_pages(temp_pdf_path, file_name)
            
            # PDF 경로를 각 Document 객체의 source 속성에 저장
            #for doc in documents:
            #    doc.source = temp_pdf_path
            
            logging.info(f"변환 완료: {temp_pdf_path}")
            return documents
        
        except Exception as e:
            logging.error(f"변환 중 오류 발생: {str(e)}")
            return None
        
        finally:
            hwp.Quit()
            pythoncom.CoUninitialize()
        
           
            # 임시 PDF 파일 삭제
            if hwp_path and os.path.exists(hwp_path):
                try:
                    os.remove(hwp_path)
                    logging.info(f"임시 hwp 파일 삭제 완료: {hwp_path}")
                except Exception as e:
                    logging.error(f"임시 hwp 파일 삭제 중 오류 발생: {str(e)}")
           
           
            # 임시 PDF 파일 삭제
            if temp_pdf_path and os.path.exists(temp_pdf_path):
                try:
                    os.remove(temp_pdf_path)
                    logging.info(f"임시 PDF 파일 삭제 완료: {temp_pdf_path}")
                except Exception as e:
                    logging.error(f"임시 PDF 파일 삭제 중 오류 발생: {str(e)}")

    def close_hwp_popups(self):
        def callback(hwnd, _):
            if win32gui.GetWindowText(hwnd) == "한글":
                win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
        
        win32gui.EnumWindows(callback, None)
    
    def handle_completion_dialog(self):
        def callback(hwnd, _):
            if "완료" in win32gui.GetWindowText(hwnd):
                # '확인' 버튼 클릭
                ok_button = win32gui.FindWindowEx(hwnd, 0, "Button", "확인")
                win32gui.PostMessage(ok_button, win32con.WM_LBUTTONDOWN, 0, 0)
                win32gui.PostMessage(ok_button, win32con.WM_LBUTTONUP, 0, 0)
        
        win32gui.EnumWindows(callback, None)
        

    def extract_text_from_pdf(self, file):
        text = ""
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")
            # 파일 객체가 읽기 모드로 열려 있는지 확인
        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")
    
        # 파일 경로인 경우 파일 열기
        if isinstance(file, str):
            file = open(file, 'rb')
        
        try:
            reader = PyPDF2.PdfReader(file)
            documents = []
            for page_num, page in enumerate(reader.pages,1):
                text += page.extract_text()
                text = re.sub(r'[^a-zA-Z0-9가-힣\s\n]', '', text)
                metadata = {
                    "source": file,
                    "file_name": file,
                    "page": page_num
                }
                doc = Document(page_content=text, metadata=metadata)
                documents.append(doc)
        finally:
            # 우리가 파일을 열었다면 닫아줍니다
            if isinstance(file, io.IOBase) and not isinstance(file, io.BytesIO):
                file.close()        
        return documents
    
    def get_file_name(self, file_object):
        name = getattr(file_object, 'name', None)
        if name:
            return os.path.basename(name)
        return "Unknown"
    
    def extract_text_from_markdown(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from Markdown file.

        Args:
            file (Union[str, IO]): Path to the Markdown file or a file-like object

        Returns:
            List[Document]: List containing a single Document object with extracted text
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                file_obj = open(file, 'r', encoding='utf-8')
            else:
                file_obj = file
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            content = file_obj.read()
            html = markdown.markdown(content)
            text = self.extract_from_html(html)

            metadata = {
                "source": file_path,
                "file_name": file_name,
                "file_type": "markdown"
            }
            doc = Document(page_content=text, metadata=metadata)
            documents.append(doc)

            return documents

        except Exception as e:
            print(f"Error processing Markdown file: {str(e)}")
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.StringIO):
                file_obj.close()

    def extract_text_from_html(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from HTML file.

        Args:
            file (Union[str, IO]): Path to the HTML file or a file-like object

        Returns:
            List[Document]: List containing a single Document object with extracted text
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                file_obj = open(file, 'r', encoding='utf-8')
            else:
                file_obj = file
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            content = file_obj.read()
            text = self.extract_from_html(content)

            metadata = {
                "source": file_path,
                "file_name": file_name,
                "file_type": "html"
            }
            doc = Document(page_content=text, metadata=metadata)
            documents.append(doc)

            return documents

        except Exception as e:
            print(f"Error processing HTML file: {str(e)}")
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.StringIO):
                file_obj.close()

    def extract_text_from_txt(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from TXT file.

        Args:
            file (Union[str, IO]): Path to the TXT file or a file-like object

        Returns:
            List[Document]: List containing a single Document object with extracted text
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                file_obj = open(file, 'r', encoding='utf-8')
            else:
                file_obj = file
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            content = file_obj.read()
            text = content.strip()

            metadata = {
                "source": file_path,
                "file_name": file_name,
                "file_type": "txt"
            }
            doc = Document(page_content=text, metadata=metadata)
            documents.append(doc)

            return documents

        except Exception as e:
            print(f"Error processing TXT file: {str(e)}")
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.StringIO):
                file_obj.close()

    def extract_from_html(self, content: str) -> str:
        """Extract text from HTML content."""
        soup = BeautifulSoup(content, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        return text